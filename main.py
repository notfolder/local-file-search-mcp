from mcp.server.fastmcp import FastMCP
import logging
from typing import Optional
from datetime import datetime
import os
from Foundation import (
    NSMetadataQuery,
    NSPredicate,
    NSRunLoop,
    NSDate,
)
import time
import platform
# オフィスファイル用のライブラリ
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

system = platform.system()
if system == 'Windows':
    import win32com.client

mcp = FastMCP("local-file-search")

# ロガーの設定
logger = logging.getLogger(__name__)
# logger.setLevel(logging.DEBUG)
handler = logging.StreamHandler()
handler.setFormatter(logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s'))
logger.addHandler(handler)


@mcp.tool()
def search_local_files(
    query: str,
    extension: Optional[str] = None,
    modified_after: Optional[str] = None,
    min_size_kb: Optional[int] = None,
    max_size_kb: Optional[int] = None
) -> str:
    """
    This function allows you to search for files in the local file system.
    Search for files in the system index based on content and various criteria.
    for files containing specific text content. The search can be refined using multiple filtering options.

    Args:
        query (str): Text to search for within file contents
        extension (Optional[str]): Filter by file extension without dot (e.g., "pdf", "txt")
        modified_after (Optional[str]): Filter files modified after this date (ISO 8601 format, e.g., "2024-01-01T00:00:00")
        min_size_kb (Optional[int]): Minimum file size in kilobytes
        max_size_kb (Optional[int]): Maximum file size in kilobytes

    Returns:
        str: A newline-separated list of found files in format "filename - path",
             or "No matching files found." if no results

    Examples:
        >>> search_local_files("report", extension="pdf", modified_after="2024-01-01T00:00:00", min_size_kb=100)
        "report.pdf - /Users/username/Documents/report.pdf"
    """
    
    if system == 'Windows':
        return search_local_files_windows(query, extension, modified_after, min_size_kb, max_size_kb)
    elif system == 'Darwin':  # macOS
        return search_local_files_mac(query, extension, modified_after, min_size_kb, max_size_kb)
    else:
        return "Unsupported operating system"

def search_local_files_mac(
    query: str,
    extension: Optional[str] = None,
    modified_after: Optional[str] = None,
    min_size_kb: Optional[int] = None,
    max_size_kb: Optional[int] = None
) -> str:
    # Documents フォルダのパスを取得
    documents_path = os.path.expanduser('~/Documents')
    logger.debug(f"Documents path: {documents_path}")
    
    # クエリの構築
    predicate_parts = [
        f'kMDItemTextContent CONTAINS[cd] "{query}"'
    ]
    
    logger.debug(f"Predicate parts: {predicate_parts}")
    
    if extension:
        predicate_parts.append(f'kMDItemFSName LIKE[cd] "*.{extension}"')

    if modified_after:
        try:
            dt = datetime.fromisoformat(modified_after)
            timestamp = time.mktime(dt.timetuple())
            date = NSDate.dateWithTimeIntervalSince1970_(timestamp)
            predicate_parts.append(f'kMDItemFSContentChangeDate >= $time')
        except ValueError:
            return "Invalid date format. Use ISO 8601 (e.g., 2024-01-01T00:00:00)"

    # NSMetadataQueryの設定
    mdquery = NSMetadataQuery.alloc().init()
    mdquery.setSearchScopes_([documents_path])  # 検索スコープを設定
    predicate = NSPredicate.predicateWithFormat_(' AND '.join(predicate_parts))
    logger.debug(f"Final predicate: {predicate}")
    mdquery.setPredicate_(predicate)
    
    # 検索開始
    mdquery.startQuery()
    
    # 結果が得られるまで待機
    NSRunLoop.currentRunLoop().runUntilDate_(NSDate.dateWithTimeIntervalSinceNow_(3.0))
    mdquery.stopQuery()
    
    # 結果の取得とフィルタリング
    filtered_files = []
    for item in mdquery.results():
        path = item.valueForAttribute_('kMDItemPath')
        if not path:
            continue
            
        try:
            size_kb = os.path.getsize(path) / 1024
            if min_size_kb and size_kb < min_size_kb:
                continue
            if max_size_kb and size_kb > max_size_kb:
                continue
                
            name = os.path.basename(path)
            kind = item.valueForAttribute_('kMDItemKind') or "Unknown"
            filtered_files.append(f"{name} ({kind}) - {path}")
        except OSError:
            continue
    logger.debug(f"Filtered files: {filtered_files}")
    
    return '\n'.join(filtered_files) if filtered_files else "No matching files found."


def search_local_files_windows(
    query: str,
    extension: Optional[str] = None,
    modified_after: Optional[str] = None,
    min_size_kb: Optional[int] = None,
    max_size_kb: Optional[int] = None
) -> str:
    """Search indexed files on Windows using Windows Search. Optionally filter by file extension, modified date, and file size range (KB)."""
    conn = win32com.client.Dispatch("ADODB.Connection")
    conn.Open("Provider=Search.CollatorDSO;Extended Properties='Application=Windows';")

    conditions = [f"CONTAINS('{query}')"]
    base_path = 'C:/Users/' + win32com.client.Dispatch("WScript.Network").UserName + '/Documents'
    conditions.append(f"SCOPE='file:{base_path}'")

    if extension:
        conditions.append(f"System.FileExtension = '{extension}'")

    if modified_after:
        try:
            dt = datetime.fromisoformat(modified_after)
            iso_time = dt.strftime('%Y-%m-%dT%H:%M:%S')
            conditions.append(f"System.DateModified >= '{iso_time}'")
        except ValueError:
            return "Invalid date format. Use ISO 8601 (e.g., 2024-01-01T00:00:00)"

    if min_size_kb is not None:
        conditions.append(f"System.Size >= {min_size_kb * 1024}")
    if max_size_kb is not None:
        conditions.append(f"System.Size <= {max_size_kb * 1024}")

    where_clause = " AND ".join(conditions)

    rs = conn.Execute(
        f"SELECT System.ItemName, System.ItemUrl FROM SYSTEMINDEX WHERE {where_clause}"
    )[0]

    results = []
    while not rs.EOF:
        name = rs.Fields.Item("System.ItemName").Value
        path = rs.Fields.Item("System.ItemUrl").Value
        results.append(f"{name} - {path}")
        rs.MoveNext()

    rs.Close()
    conn.Close()

    return "\n".join(results) if results else "No matching files found."


def read_word_file(file_path: str) -> str:
    """Wordファイルをpython-docxで読み取る"""
    try:
        doc = Document(file_path)
        return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    except ModuleNotFoundError:
        return "[ERROR] python-docx is not installed. Run: pip install python-docx"
    except Exception as e:
        return f"[ERROR] Failed to read Word file: {e}"

def read_excel_file(file_path: str) -> str:
    """Excelファイルをopenpyxlで読み取る"""
    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.rows:
                row_values = [str(cell.value) if cell.value is not None else '' for cell in row]
                text.append('\t'.join(row_values))
        return '\n'.join(text)
    except ModuleNotFoundError:
        return "[ERROR] openpyxl is not installed. Run: pip install openpyxl"
    except Exception as e:
        return f"[ERROR] Failed to read Excel file: {e}"

def read_ppt_file(file_path: str) -> str:
    """PowerPointファイルをpython-pptxで読み取る"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except ModuleNotFoundError:
        return "[ERROR] python-pptx is not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"[ERROR] Failed to read PowerPoint file: {e}"

@mcp.tool()
def local_read_file(file_path: str)-> str:
    """Read and extract text content from various file types including Office documents.

    This function attempts to read the content of a file, first trying as a text file,
    then using specialized libraries for Office documents (Word, Excel, PowerPoint).

    Args:
        file_path (str): Full path to the file to be read

    Returns:
        str: Text content of the file if successful,
             Error message starting with '[ERROR]' if reading fails,
             '[SKIP]' message for unsupported file types

    Examples:
        >>> local_read_file("/path/to/document.docx")
        "Content of the Word document..."
        >>> local_read_file("/path/to/unsupported.xyz")
        "[SKIP] Unsupported file type: .xyz"
    """
    ext = os.path.splitext(file_path)[1].lower()

    # まずテキストファイルとして試みる
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except:
        pass

    # Pythonライブラリで読み取りを試行
    try:
        if ext in [".doc", ".docx"]:
            return read_word_file(file_path)
        elif ext in [".xls", ".xlsx"]:
            return read_excel_file(file_path)
        elif ext in [".ppt", ".pptx"]:
            return read_ppt_file(file_path)
        else:
            return f"[SKIP] Unsupported file type: {ext}"
    except Exception as e:
        return f"[ERROR] Failed to read {file_path}: {e}"


if __name__ == "__main__":
    mcp.run()
    # mcp.run(transport="sse")
